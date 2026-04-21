#!/usr/bin/env python3
"""
北信科 PPT 一键生成
基于 python-pptx，使用北信科官方 PPT 模板生成演示文稿
"""

import argparse
import io
import json
import os
import re
import sys
from pathlib import Path
from typing import Optional

# Windows GBK 编码兼容
if sys.platform == "win32":
    try:
        if hasattr(sys.stdout, 'buffer') and sys.stdout.encoding != 'utf-8':
            sys.stdout.reconfigure(encoding='utf-8', errors='replace')
        if hasattr(sys.stderr, 'buffer') and sys.stderr.encoding != 'utf-8':
            sys.stderr.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        pass

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("❌ 需要安装 python-pptx: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


SCRIPT_DIR = Path(__file__).parent
PROJECT_DIR = SCRIPT_DIR.parent
TEMPLATES_DIR = PROJECT_DIR / "templates"


def list_templates():
    """列出可用的 PPT 模板"""
    if not TEMPLATES_DIR.exists():
        print(f"📭 模板目录不存在: {TEMPLATES_DIR}")
        print("   请将 PPT 模板文件放入 templates/ 目录")
        return

    templates = list(TEMPLATES_DIR.glob("*.pptx"))
    if not templates:
        print("📭 暂无可用模板")
        return

    print("🎨 可用 PPT 模板：")
    for t in templates:
        size_kb = t.stat().st_size / 1024
        print(f"  • {t.name} ({size_kb:.1f} KB)")


def generate_ppt(
    title: str,
    markdown_file: str = "",
    content: str = "",
    template: str = "",
    output: str = "output.pptx",
    no_polish: bool = False,
):
    """
    生成 PPT

    Args:
        title: 演示文稿标题
        markdown_file: Markdown 内容文件路径
        content: 直接提供的内容文本
        template: 模板文件名
        output: 输出文件路径
        no_polish: 是否关闭文字优化
    """
    # 读取内容
    if markdown_file:
        md_path = Path(markdown_file)
        if not md_path.exists():
            print(f"❌ Markdown 文件不存在: {md_path}", file=sys.stderr)
            sys.exit(1)
        content = md_path.read_text(encoding="utf-8")

    if not content:
        content = f"# {title}\n\n请在 SKILL.md 指引下使用 AI 生成内容"

    # 解析 Markdown 为幻灯片
    slides_data = parse_markdown(content, title)

    # 加载模板或创建空白演示
    if template:
        template_path = TEMPLATES_DIR / template
        if template_path.exists():
            prs = Presentation(str(template_path))
            # 清空模板中的示例页
            while len(prs.slides) > 0:
                rId = prs.slides._sldIdLst[0].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[0]
        else:
            print(f"⚠️ 模板不存在: {template_path}，使用空白模板", file=sys.stderr)
            prs = Presentation()
    else:
        prs = Presentation()

    # 添加幻灯片
    for slide_data in slides_data:
        add_slide(prs, slide_data, no_polish)

    # 保存
    prs.save(output)
    print(f"✅ PPT 已生成: {output}")
    print(f"   标题: {title}")
    print(f"   页数: {len(slides_data)}")


def parse_markdown(content: str, default_title: str = "") -> list:
    """
    将 Markdown 内容解析为幻灯片数据

    每个一级标题 (# ) 作为新幻灯片的标题
    二级标题 (## ) 作为内容分块
    """
    slides = []
    current_slide = None

    for line in content.split("\n"):
        if line.startswith("# ") and not line.startswith("## "):
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                "title": line[2:].strip(),
                "content": [],
            }
        elif line.startswith("## "):
            if current_slide is None:
                current_slide = {"title": default_title, "content": []}
            current_slide["content"].append({"type": "heading", "text": line[3:].strip()})
        elif line.strip():
            if current_slide is None:
                current_slide = {"title": default_title, "content": []}
            current_slide["content"].append({"type": "text", "text": line.strip()})

    if current_slide:
        slides.append(current_slide)

    if not slides:
        slides.append({"title": default_title, "content": [{"type": "text", "text": content.strip()}]})

    return slides


def add_slide(prs: Presentation, slide_data: dict, no_polish: bool = False):
    """添加一页幻灯片"""
    # 使用空白布局
    blank_layout = None
    for layout in prs.slide_layouts:
        if "blank" in layout.name.lower() or layout.name == "空白":
            blank_layout = layout
            break
    if blank_layout is None:
        blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]

    slide = prs.slides.add_slide(blank_layout)

    # 添加标题
    title = slide_data.get("title", "")
    if title:
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.5), Inches(8.4), Inches(1.0)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT

        # 北信科主题色 - 蓝色标题下划线
        from pptx.dml.color import RGBColor
        p.font.color.rgb = RGBColor(0, 70, 140)  # BISTU 蓝

    # 添加内容
    content_items = slide_data.get("content", [])
    if content_items:
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        for i, item in enumerate(content_items):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            text = item.get("text", "")
            if item.get("type") == "heading":
                p.text = text
                p.font.size = Pt(22)
                p.font.bold = True
                p.space_before = Pt(12)
            else:
                # 处理列表项
                if text.startswith("- ") or text.startswith("* "):
                    text = "• " + text[2:]
                elif re.match(r"^\d+\.\s", text):
                    pass  # 保持有序列表
                p.text = text
                p.font.size = Pt(18)
                p.space_before = Pt(6)

            p.alignment = PP_ALIGN.LEFT

    # 添加页脚
    footer_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(6.8), Inches(8.4), Inches(0.4)
    )
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "北京信息科技大学 | BISTU"
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(128, 128, 128)
    p.alignment = PP_ALIGN.RIGHT


def main():
    parser = argparse.ArgumentParser(description="北信科 PPT 生成器")
    parser.add_argument("--list-templates", action="store_true", help="列出可用模板")
    parser.add_argument("--title", type=str, help="PPT 标题")
    parser.add_argument("--markdown", type=str, help="Markdown 内容文件路径")
    parser.add_argument("--content", type=str, help="直接提供的内容文本")
    parser.add_argument("--template", type=str, default="", help="模板文件名")
    parser.add_argument("--output", type=str, default="output.pptx", help="输出文件路径")
    parser.add_argument("--no-polish", action="store_true", help="关闭文字优化")

    args = parser.parse_args()

    if args.list_templates:
        list_templates()
        return

    if not args.title:
        print("❌ 请提供 --title 参数", file=sys.stderr)
        parser.print_help()
        return

    generate_ppt(
        title=args.title,
        markdown_file=args.markdown,
        content=args.content,
        template=args.template,
        output=args.output,
        no_polish=args.no_polish,
    )


if __name__ == "__main__":
    main()
